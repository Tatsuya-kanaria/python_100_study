# %%
import pptx
import docx


# pptx ----------------------------------

pptx_data = pptx.Presentation('./data/サンプル_PowerPoint.pptx')
# スライド数
len(pptx_data.slides)

sld_0 = pptx_data.slides[0]
shp_sld_0 = sld_0.shapes
# シェープの数
len(shp_sld_0)

print(shp_sld_0[0].text)
# 中身がテキストならTrue
print(shp_sld_0[0].has_text_frame)

# 全文の取得
texts = []
for slide in pptx_data.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text)
print(f'pptx_texts: ', texts)
print()

# docx ----------------------------------

docx_data = docx.Document('./data/サンプル_Word.docx')
# 改行の数
len(docx_data.paragraphs)

docx_data.paragraphs[0].text

texts = []
for paragraph in docx_data.paragraphs:
    texts.append(paragraph.text)
print(f'docx_texts: ', texts)

# %%
